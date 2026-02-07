import cv2
import numpy as np
import os
from pptx import Presentation
from pptx.util import Inches


class ImageEngine:
    def __init__(self):
        self.asset_dir = os.path.join("assets", "images")
        os.makedirs(self.asset_dir, exist_ok=True)

    def get_image_path(self, char):
        """Returns path for a character image (A-Z)."""
        char = char.lower()
        for ext in [".png", ".jpg", ".jpeg"]:
            path = os.path.join(self.asset_dir, char + ext)
            if os.path.exists(path):
                return path
        return None

    def create_word_strip(self, word):
        """Creates a horizontal image strip of the finger-spelled word."""
        images = []
        for char in word:
            if char.isalpha():
                p = self.get_image_path(char)
                if p:
                    img = cv2.imread(p)
                    if img is not None:
                        # Resize to fixed height 100px
                        h, w, _ = img.shape
                        scale = 100 / h
                        img = cv2.resize(img, (int(w * scale), 100))
                        images.append(img)

        if not images: return None
        return np.hstack(images)

    def generate_ppt(self, text, output_path):
        """Generates a PowerPoint presentation."""
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "ASL Translation"
        title_slide.placeholders[1].text = text

        words = text.split()
        for word in words:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = word.upper()

            strip = self.create_word_strip(word)
            if strip is not None:
                t_path = "temp_strip.png"
                cv2.imwrite(t_path, strip)
                # Add picture centered
                slide.shapes.add_picture(t_path, Inches(1), Inches(2), width=Inches(8))
                os.remove(t_path)

        prs.save(output_path)
        return output_path